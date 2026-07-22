import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

def generate_expanded_powerpoint():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Color Palette
    DARK_BG = RGBColor(11, 15, 23)
    CYAN_ACCENT = RGBColor(6, 182, 212)
    PURPLE_ACCENT = RGBColor(139, 92, 246)
    AMBER_ACCENT = RGBColor(245, 158, 11)
    ROSE_ACCENT = RGBColor(244, 63, 94)
    WHITE_TEXT = RGBColor(248, 250, 252)
    SLATE_TEXT = RGBColor(148, 163, 184)

    slides_data = [
        {
            "title": "AURA-8: Unified Asset & Operations Brain",
            "subtitle": "Pitch Presentation Deck | Economic Times AI Hackathon 2.0 (2026)",
            "bullets": [
                "Problem Statement #8: AI for Industrial Knowledge Intelligence: Unified Asset & Operations Brain",
                "Theme: Industrial Intelligence / Document Management / Knowledge Engineering / Quality",
                "Core Innovation: Universal Document Ingestion, Dynamic 2D Knowledge Graph, and Grounded Gemini RAG Copilot",
                "Engineered Stack: Next.js 15 App Router, React 19, Python 3.11 FastAPI, pdf-parse, and Google Gemini API (gemini-flash-latest)"
            ]
        },
        {
            "title": "Slide 2: The Problem Context in Heavy Industry",
            "subtitle": "Why Indian Manufacturing & Refining Facilities Lose Millions to Disconnected Data",
            "bullets": [
                "35% Lost Working Hours: Industrial engineers and operators spend over a third of their shift searching for drawings, clarifying SOPs, or recreating lost records.",
                "7 to 12 Disconnected Document Systems: Plants operate across P&ID CAD drawings in one place, maintenance work orders in another, operating procedures in a third, and regulatory filings in email archives.",
                "18–22% Unplanned Downtime: Maintenance teams make critical repair decisions without complete equipment history or failure pattern context.",
                "The Knowledge Cliff: 25% of India's experienced industrial engineers will retire within the next decade, taking decades of undocumented operational knowledge with them."
            ]
        },
        {
            "title": "Slide 3: The Solution — AURA-8 Platform Overview",
            "subtitle": "A Single, Unified Neural Brain for Industrial Plant Knowledge",
            "bullets": [
                "Universal Document Ingestion Engine: Ingests heterogeneous engineering drawings (P&IDs), OEM manuals, thermography logs, work orders, and statutory standards across PDF, TXT, and CSV formats.",
                "Dynamic 2D Visual Knowledge Graph: Automatically parses equipment tags (P-101A, V-301), parameters, and OISD/Factory Act rules into an interactive visual HTML5 canvas network.",
                "Persistent Gemini RAG Copilot: Floating AI assistant drawer accessible across all pages, answering operational queries with exact page-level text citations.",
                "Statutory Audit Automation: 1-Click export of downloadable compliance evidence packages for OISD-STD-137, Factory Act 1948 Section 37, and PESO rules."
            ]
        },
        {
            "title": "Slide 4: Gateway Landing & Dual-Mode Operating Model",
            "subtitle": "Clean Workspace vs Guest Evaluator Mode",
            "bullets": [
                "Start Clean Workspace Mode: Pristine zero-data state. Infographics, Knowledge Graph nodes, and RAG indexes populate ONLY after the user submits their first document.",
                "Guest Demo Mode: Instant 1-click access pre-loaded with sample hydrocracker refinery data (P-101A pump, V-301 column, OISD-137 compliance) for hackathon evaluators.",
                "Global Mode Switcher: One-click mode switching available on EVERY page header across the platform.",
                "User Privacy & Security: API keys managed server-side via .env.local without exposing credentials in the client UI."
            ]
        },
        {
            "title": "Slide 5: Technical Architecture & Execution Pipeline",
            "subtitle": "Next.js 15 App Router + Python FastAPI + Google Gemini API",
            "bullets": [
                "Client Layer (Next.js 15): React 19 App Router with Tailwind CSS, Lucide icons, HTML5 Canvas 2D Graph visualizer, and persistent slide-over chatbot drawer.",
                "API Orchestration (/api/upload & /api/query): Native Next.js server-side routes handling multipart uploads and vector RAG requests with zero CORS issues.",
                "Text Stream Decompression: pdf-parse decompresses PDF /FlateDecode binary streams page-by-page, converting binary code into clean human-readable text.",
                "Google Gemini RAG Engine: Calls https://generativelanguage.googleapis.com/v1beta/models/gemini-flash-latest:generateContent with X-goog-api-key for grounded synthesis."
            ]
        },
        {
            "title": "Slide 6: Dynamic 2D Industrial Knowledge Graph",
            "subtitle": "Automated Entity Extraction & Statutory Linkage",
            "bullets": [
                "RegEx & NLP Entity Extractor: Parses equipment tags (P-101A, V-301), process parameters (310°C, 4.9 bar), and regulatory references (OISD-137, Factory Act §37).",
                "Interactive HTML5 Canvas Graph: Renders nodes color-coded by category (Equipment, OEM Manual, Regulation, Incident, Work Order).",
                "Right-Side Asset Inspector: Displays live technical attributes (Operating Temp 84.5°C, Barrier Pressure 4.9 bar), risk status (High Warning), and direct button to query AI Copilot on that asset.",
                "Graph Expansion Velocity: Tracks monthly ingested documents and extracted semantic entities dynamically."
            ]
        },
        {
            "title": "Slide 7: Expert Gemini RAG Copilot & Persistent Chatbot",
            "subtitle": "Grounded Conversational AI with Direct Page Citations",
            "bullets": [
                "Persistent Floating Chatbot: Slide-over drawer anchored on every page, allowing operators to query plant documents without losing context.",
                "Vector Cosine Match Scoring: Calculates exact term-frequency cosine similarity score for retrieved text chunks.",
                "Grounded Page Citations: Clickable citation chips open a side-by-side drawer showing the exact text snippet extracted from the PDF page.",
                "Field-Technician Ready: Built with mobile-responsive design for field engineers using tablets and rugged smartphones."
            ]
        },
        {
            "title": "Slide 8: Predictive Maintenance & Root Cause Analysis (RCA)",
            "subtitle": "Ishikawa Fishbone Synthesis & Remaining Useful Life (RUL) Trajectory",
            "bullets": [
                "Remaining Useful Life (RUL) Forecast: Models degradation trajectories (38.5 Hours Until Trip Limit for P-101A) based on vibration FFT telemetry (2.7 mm/s RMS) and bearing temp (84.5°C).",
                "Automated Ishikawa (Fishbone) RCA: Fuses work order logs, OEM tolerances, and near-miss incident history into a 4-branch Root Cause diagram (Thermal Stress, Barrier Fluid System, Material Specs, Operations).",
                "Preventive Work Order Generation: 1-Click button to auto-issue preventive work orders before an unplanned trip limit breach."
            ]
        },
        {
            "title": "Slide 9: Quality & Statutory Compliance Intelligence",
            "subtitle": "Automated Audit Monitoring & Evidence Package Export",
            "bullets": [
                "Continuous Statutory Audit: Monitors plant operations against OISD-STD-137 Clause 4.2.1 (dual mechanical seal pressurization), Factory Act 1948 Section 37 (precaution against explosive fumes), and PESO Rules 2016.",
                "Non-Compliance Gap Alerts: Highlights specific corrective remediation requirements (e.g. repressurize Plan 53B barrier fluid bladder).",
                "Downloadable Audit Evidence Package: 1-Click export of downloadable compliance report files (OISD_FactoryAct_Compliance_Package_Q3_2026.txt)."
            ]
        },
        {
            "title": "Slide 10: Value Proposition, Business Impact & Hackathon Jury Alignment",
            "subtitle": "Demonstrated ROI & Scoring Criteria Alignment",
            "bullets": [
                "Innovation (25%): Real-time conversion of unstructured PDFs into visual 2D Knowledge Graphs & Gemini-grounded RAG answers.",
                "Business Impact (25%): 18-22% reduction in unplanned downtime events and elimination of fatal industrial safety non-compliance.",
                "Technical Excellence (20%): Decoupled Next.js 15 + Python FastAPI + pdf-parse + Google Gemini API architecture.",
                "User Experience (15%): Gateway landing page, zero-data initial state, global mode switcher, and persistent chatbot across all views."
            ]
        }
    ]

    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Solid Dark Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        # Top Banner Accent Bar
        accent_bar = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(0), Inches(13.333), Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = CYAN_ACCENT
        accent_bar.line.color.rgb = CYAN_ACCENT

        # Title Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.3))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = slide_info["title"]
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT

        p2 = tf.add_paragraph()
        p2.text = slide_info["subtitle"]
        p2.font.size = Pt(13)
        p2.font.color.rgb = SLATE_TEXT

        # Content Card Box
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.733), Inches(4.8))
        ctf = content_box.text_frame
        ctf.word_wrap = True

        for bullet in slide_info["bullets"]:
            bp = ctf.add_paragraph()
            bp.text = bullet
            bp.font.size = Pt(16)
            bp.font.color.rgb = WHITE_TEXT
            bp.space_after = Pt(14)

    output_path = r"c:\Users\sukri\Desktop\ET_AI_Hackathon_problem\AURA8_Presentation_Deck.pptx"
    prs.save(output_path)
    print(f"Expanded 10-Slide Presentation saved to {output_path}")

if __name__ == "__main__":
    generate_expanded_powerpoint()
